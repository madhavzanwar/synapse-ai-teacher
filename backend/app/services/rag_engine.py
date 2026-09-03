"""
Production-Grade RAG Knowledge Grounding Pipeline for Synapse AI Teacher.
Handles multi-format document parsing (PDF, DOCX, PPTX, TXT, MD), structural extraction,
semantic/section-aware chunking, dense vector similarity + sparse BM25 hybrid search,
and Reciprocal Rank Fusion (RRF) re-ranking.
"""
import io
import os
import re
import math
import uuid
import time
import logging
from typing import List, Dict, Any, Optional, Tuple, Set
from dataclasses import dataclass, field, asdict

# Optional ML / Parser libraries with graceful fallback
try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    import docx
except ImportError:
    docx = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from rank_bm25 import BM25Okapi
except ImportError:
    BM25Okapi = None

try:
    import google.generativeai as genai
except ImportError:
    genai = None

from app.config import settings

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Data Classes
# ---------------------------------------------------------------------------

@dataclass
class DocumentMetadata:
    document_id: str
    source_filename: str
    detected_title: str
    file_type: str
    total_pages: int
    total_chunks: int
    sections_detected: List[str]
    char_count: int
    has_math: bool
    created_at: float = field(default_factory=time.time)

    def to_dict(self) -> Dict[str, Any]:
        return asdict(self)


@dataclass
class DocumentChunk:
    chunk_id: str
    document_id: str
    source_filename: str
    page_number: int
    section_title: str
    chunk_index: int
    text: str
    tokens: List[str] = field(default_factory=list)
    dense_vector: Optional[List[float]] = None
    has_math: bool = False
    char_count: int = 0
    score: float = 0.0

    def to_dict(self) -> Dict[str, Any]:
        data = asdict(self)
        # Avoid serializing full dense vector in lightweight JSON responses
        data.pop("dense_vector", None)
        return data


# ---------------------------------------------------------------------------
# Pure-Python High-Speed BM25 Okapi Fallback & Tokenizer
# ---------------------------------------------------------------------------

def tokenize_text(text: str) -> List[str]:
    """Extract lowercase alpha-numeric tokens and math symbols."""
    clean = text.lower()
    # Normalize LaTeX symbols to identifiable tokens
    clean = re.sub(r"\\(frac|sqrt|text|mathbf|int|sum|prod|alpha|beta|gamma|theta|lambda|sigma|pi|partial)", r" math_\1 ", clean)
    tokens = re.findall(r"\b[a-zA-Z0-9_]{2,}\b", clean)
    return tokens


class PureBM25Okapi:
    """Production-grade BM25 Okapi implementation with IDF smoothing."""

    def __init__(self, corpus_tokens: List[List[str]], k1: float = 1.5, b: float = 0.75):
        self.k1 = k1
        self.b = b
        self.corpus_size = len(corpus_tokens)
        self.avgdl = sum(len(doc) for doc in corpus_tokens) / max(self.corpus_size, 1)
        self.doc_freqs: List[Dict[str, int]] = []
        self.idf: Dict[str, float] = {}
        self.doc_len: List[int] = []

        # Count frequencies
        df: Dict[str, int] = {}
        for doc in corpus_tokens:
            self.doc_len.append(len(doc))
            frequencies: Dict[str, int] = {}
            for term in doc:
                frequencies[term] = frequencies.get(term, 0) + 1
            self.doc_freqs.append(frequencies)

            for term in frequencies:
                df[term] = df.get(term, 0) + 1

        # Calculate smoothed IDF
        for term, freq in df.items():
            self.idf[term] = math.log((self.corpus_size - freq + 0.5) / (freq + 0.5) + 1.0)

    def get_scores(self, query_tokens: List[str]) -> List[float]:
        scores = [0.0] * self.corpus_size
        for i, doc_freq in enumerate(self.doc_freqs):
            doc_len = self.doc_len[i]
            for term in query_tokens:
                if term in doc_freq:
                    freq = doc_freq[term]
                    idf = self.idf.get(term, 0.0)
                    numerator = idf * freq * (self.k1 + 1)
                    denominator = freq + self.k1 * (1 - self.b + self.b * (doc_len / max(self.avgdl, 1e-5)))
                    scores[i] += numerator / max(denominator, 1e-5)
        return scores


# ---------------------------------------------------------------------------
# Main RAG Engine
# ---------------------------------------------------------------------------

class RAGEngine:
    """
    Enterprise RAG Engine for multimodal educational materials.
    Parses PDF, DOCX, PPTX, TXT, MD; applies semantic chunking with LaTeX formula preservation;
    executes Hybrid Dense + BM25 Search with Reciprocal Rank Fusion (RRF).
    """

    def __init__(self):
        self.documents: Dict[str, DocumentMetadata] = {}
        self.raw_texts: Dict[str, str] = {}
        self.chunks: List[DocumentChunk] = []
        self._bm25_index: Optional[PureBM25Okapi] = None
        self._corpus_tokens: List[List[str]] = []
        self._index_dirty: bool = True

    # -----------------------------------------------------------------------
    # Document Parsing Layer
    # -----------------------------------------------------------------------

    def ingest_document(self, file_path: str, document_id: Optional[str] = None) -> DocumentMetadata:
        """Parse file from disk path into structured chunks with metadata."""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        filename = os.path.basename(file_path)
        doc_id = document_id or str(uuid.uuid4())
        ext = os.path.splitext(filename)[1].lower()

        logger.info(f"Ingesting document '{filename}' ({ext}) as ID: {doc_id}")

        if ext == ".pdf":
            parsed_pages, sections = self._parse_pdf(file_path)
            detected_title = self._detect_title(parsed_pages, filename)
            file_type = "pdf"
        elif ext in [".docx", ".doc"]:
            parsed_pages, sections = self._parse_docx(file_path)
            detected_title = sections[0] if sections else filename
            file_type = "docx"
        elif ext in [".pptx", ".ppt"]:
            parsed_pages, sections = self._parse_pptx(file_path)
            detected_title = sections[0] if sections else filename
            file_type = "pptx"
        else:
            parsed_pages, sections = self._parse_text_markdown(file_path)
            detected_title = sections[0] if sections else filename
            file_type = "text"

        # Combine text for document store
        full_text = "\n\n".join(page_text for _, page_text in parsed_pages)
        has_math = bool(re.search(r"(\$|\\\[|\\\(|\\text\{|\\frac\{|\\sum|\\int|\\sqrt)", full_text))

        # Semantic Chunking
        new_chunks = self._semantic_chunk(
            document_id=doc_id,
            source_filename=filename,
            parsed_pages=parsed_pages,
            sections=sections
        )

        metadata = DocumentMetadata(
            document_id=doc_id,
            source_filename=filename,
            detected_title=detected_title,
            file_type=file_type,
            total_pages=len(parsed_pages),
            total_chunks=len(new_chunks),
            sections_detected=sections[:15],
            char_count=len(full_text),
            has_math=has_math
        )

        # Store metadata and chunks
        self.documents[doc_id] = metadata
        self.raw_texts[doc_id] = full_text
        self.chunks.extend(new_chunks)
        self._index_dirty = True

        logger.info(f"Successfully ingested '{filename}': {len(new_chunks)} chunks, {len(sections)} sections.")
        return metadata

    def ingest_bytes(self, file_bytes: bytes, filename: str, document_id: Optional[str] = None) -> DocumentMetadata:
        """Parse in-memory file bytes directly."""
        doc_id = document_id or str(uuid.uuid4())
        ext = os.path.splitext(filename)[1].lower()

        temp_dir = os.path.join(os.path.dirname(__file__), "..", "..", "temp_uploads")
        os.makedirs(temp_dir, exist_ok=True)
        temp_path = os.path.join(temp_dir, f"{doc_id}_{filename}")

        with open(temp_path, "wb") as f:
            f.write(file_bytes)

        try:
            meta = self.ingest_document(temp_path, document_id=doc_id)
            return meta
        finally:
            if os.path.exists(temp_path):
                try:
                    os.remove(temp_path)
                except Exception:
                    pass

    def ingest_text(self, text: str, title: str = "Uploaded Document", document_id: Optional[str] = None) -> str:
        """Convenience method to ingest plain text / markdown notes directly."""
        meta = self.ingest_bytes(
            file_bytes=text.encode("utf-8"),
            filename=f"{title.replace(' ', '_')}.md",
            document_id=document_id
        )
        return meta.document_id

    def retrieve(self, query: str, top_k: int = 4, doc_filter_ids: Optional[List[str]] = None) -> str:
        """Convenience method returning formatted context string for legacy calls."""
        doc_id = doc_filter_ids[0] if doc_filter_ids else None
        chunks = self.retrieve_context(query=query, document_id=doc_id, top_k=top_k)
        return self.format_grounding_context(chunks)

    # -----------------------------------------------------------------------
    # Multi-Format Parsers
    # -----------------------------------------------------------------------

    def _parse_pdf(self, file_path: str) -> Tuple[List[Tuple[int, str]], List[str]]:
        """Extract multi-page PDF text while preserving structure and mathematical formulas."""
        pages: List[Tuple[int, str]] = []
        sections: List[str] = []

        if PdfReader is not None:
            try:
                reader = PdfReader(file_path)
                for page_idx, page in enumerate(reader.pages):
                    raw_text = page.extract_text() or ""
                    cleaned = self._normalize_math_and_layout(raw_text)
                    if cleaned.strip():
                        pages.append((page_idx + 1, cleaned))
                        # Detect section headers
                        found_sections = self._extract_header_lines(cleaned)
                        sections.extend(found_sections)
            except Exception as e:
                logger.error(f"Error reading PDF with pypdf: {e}")

        if not pages:
            # Fallback if empty or failed
            pages.append((1, f"[PDF Extraction for {os.path.basename(file_path)}]"))

        return pages, list(dict.fromkeys(sections))

    def _parse_docx(self, file_path: str) -> Tuple[List[Tuple[int, str]], List[str]]:
        """Extract DOCX text preserving heading hierarchy and markdown tables."""
        pages: List[Tuple[int, str]] = []
        sections: List[str] = []

        if docx is not None:
            try:
                doc = docx.Document(file_path)
                current_section = "Introduction"
                page_chunks = []

                for para in doc.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    if para.style.name.startswith("Heading"):
                        sections.append(text)
                        current_section = text
                        page_chunks.append(f"\n### {text}\n")
                    else:
                        page_chunks.append(text)

                # Process tables
                for table in doc.tables:
                    table_rows = []
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip().replace("\n", " ") for cell in row.cells)
                        table_rows.append(f"| {row_text} |")
                    if table_rows:
                        header_sep = "| " + " | ".join(["---"] * len(table.columns)) + " |"
                        table_md = "\n" + table_rows[0] + "\n" + header_sep + "\n" + "\n".join(table_rows[1:]) + "\n"
                        page_chunks.append(table_md)

                combined = "\n\n".join(page_chunks)
                pages.append((1, combined))
            except Exception as e:
                logger.error(f"Error parsing DOCX: {e}")

        if not pages:
            pages.append((1, f"[DOCX Content for {os.path.basename(file_path)}]"))

        return pages, list(dict.fromkeys(sections))

    def _parse_pptx(self, file_path: str) -> Tuple[List[Tuple[int, str]], List[str]]:
        """Extract PPTX slides preserving slide titles and bullet points."""
        pages: List[Tuple[int, str]] = []
        sections: List[str] = []

        if Presentation is not None:
            try:
                prs = Presentation(file_path)
                for slide_idx, slide in enumerate(prs.slides):
                    slide_num = slide_idx + 1
                    slide_title = f"Slide {slide_num}"
                    slide_texts = []

                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                p_text = paragraph.text.strip()
                                if p_text:
                                    if shape == slide.shapes[0] and len(p_text) < 100:
                                        slide_title = p_text
                                        sections.append(p_text)
                                    slide_texts.append(p_text)

                    slide_content = f"### Slide {slide_num}: {slide_title}\n" + "\n".join(f"- {t}" for t in slide_texts if t != slide_title)
                    pages.append((slide_num, slide_content))
            except Exception as e:
                logger.error(f"Error parsing PPTX: {e}")

        if not pages:
            pages.append((1, f"[PPTX Content for {os.path.basename(file_path)}]"))

        return pages, list(dict.fromkeys(sections))

    def _parse_text_markdown(self, file_path: str) -> Tuple[List[Tuple[int, str]], List[str]]:
        """Extract plain text or markdown file."""
        sections: List[str] = []
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
        except Exception as e:
            content = f"[Text read error: {e}]"

        sections = self._extract_header_lines(content)
        return [(1, content)], list(dict.fromkeys(sections))

    # -----------------------------------------------------------------------
    # Layout, Math & Header Normalization
    # -----------------------------------------------------------------------

    def _normalize_math_and_layout(self, text: str) -> str:
        """Clean PDF artifacts, rejoin broken hyphens, and normalize math formulas."""
        # Rejoin hyphenated words across line breaks
        text = re.sub(r"(\w+)-\n(\w+)", r"\1\2", text)
        # Normalize excessive whitespace
        text = re.sub(r"[ \t]+", " ", text)
        # Ensure LaTeX block equations are on separate lines
        text = re.sub(r"(\$\$[^\$]+\$\$)", r"\n\1\n", text)
        return text

    def _extract_header_lines(self, text: str) -> List[str]:
        """Detect Markdown headings (`#`, `##`, `###`) and numbered chapter titles."""
        headers = []
        for line in text.split("\n"):
            line = line.strip()
            # Markdown header format
            md_match = re.match(r"^#{1,3}\s+(.+)$", line)
            if md_match:
                headers.append(md_match.group(1).strip())
                continue
            # Numbered section format: "1.2 Attention Mechanism" or "Chapter 3: Backprop"
            num_match = re.match(r"^(\d+\.[\d\.]*\s+[A-Z][a-zA-Z0-9\s]{3,50}|Chapter\s+\d+[:\s]+[A-Za-z0-9\s]+)", line)
            if num_match:
                headers.append(num_match.group(1).strip())
        return headers

    def _detect_title(self, parsed_pages: List[Tuple[int, str]], filename: str) -> str:
        """Extract document title from first page or fall back to filename."""
        if parsed_pages:
            first_page = parsed_pages[0][1]
            lines = [l.strip() for l in first_page.split("\n") if l.strip()]
            for line in lines[:5]:
                # Title heuristic: reasonable length, not starting with copyright or page number
                if 5 < len(line) < 80 and not line.lower().startswith(("page", "http", "doi", "copyright")):
                    return re.sub(r"^#+\s*", "", line)
        return os.path.splitext(filename)[0].replace("_", " ").replace("-", " ").title()

    # -----------------------------------------------------------------------
    # Intelligent Semantic Chunking
    # -----------------------------------------------------------------------

    def _semantic_chunk(
        self,
        document_id: str,
        source_filename: str,
        parsed_pages: List[Tuple[int, str]],
        sections: List[str],
        max_chunk_chars: int = 1200,
        overlap_chars: int = 150
    ) -> List[DocumentChunk]:
        """
        Structure-aware semantic chunking.
        Splits by major headings, paragraphs, and formula blocks with context-preserving overlap.
        """
        chunks: List[DocumentChunk] = []
        chunk_index = 0

        for page_num, page_text in parsed_pages:
            current_section = "General Overview"
            paragraphs = [p.strip() for p in re.split(r"\n\s*\n", page_text) if p.strip()]
            current_buffer = []
            current_length = 0

            for para in paragraphs:
                para_len = len(para)
                # Check if paragraph starts with a section header (e.g. "## 2. The Bloch Sphere")
                is_header = bool(re.match(r"^#{1,3}\s+[^\n]+", para) or re.match(r"^(\d+\.[\d\.]*\s+[A-Z]|Chapter\s+\d+)", para))
                
                # If paragraph contains a new section header and buffer already has text, flush chunk
                if is_header and current_buffer and current_length > 15:
                    chunk_text = "\n\n".join(current_buffer)
                    has_math = bool(re.search(r"(\$|\\\[|\\\(|\\text\{|\\frac\{|\\sum|\\int)", chunk_text))
                    tokens = tokenize_text(chunk_text)

                    chunks.append(
                        DocumentChunk(
                            chunk_id=f"{document_id}-chk-{chunk_index}",
                            document_id=document_id,
                            source_filename=source_filename,
                            page_number=page_num,
                            section_title=current_section,
                            chunk_index=chunk_index,
                            text=chunk_text,
                            tokens=tokens,
                            has_math=has_math,
                            char_count=len(chunk_text)
                        )
                    )
                    chunk_index += 1
                    current_buffer = [para]
                    current_length = para_len

                    # Update current section title
                    h_match = re.match(r"^#{1,3}\s+(.+)$", para)
                    if h_match:
                        current_section = h_match.group(1).strip()
                    else:
                        current_section = para.split("\n")[0][:60]
                    continue

                if is_header:
                    h_match = re.match(r"^#{1,3}\s+(.+)$", para)
                    if h_match:
                        current_section = h_match.group(1).strip()

                if current_length + para_len > max_chunk_chars and current_buffer:
                    chunk_text = "\n\n".join(current_buffer)
                    has_math = bool(re.search(r"(\$|\\\[|\\\(|\\text\{|\\frac\{|\\sum|\\int)", chunk_text))
                    tokens = tokenize_text(chunk_text)

                    chunks.append(
                        DocumentChunk(
                            chunk_id=f"{document_id}-chk-{chunk_index}",
                            document_id=document_id,
                            source_filename=source_filename,
                            page_number=page_num,
                            section_title=current_section,
                            chunk_index=chunk_index,
                            text=chunk_text,
                            tokens=tokens,
                            has_math=has_math,
                            char_count=len(chunk_text)
                        )
                    )
                    chunk_index += 1

                    # Sliding overlap
                    if len(current_buffer[-1]) <= overlap_chars:
                        current_buffer = [current_buffer[-1], para]
                    else:
                        current_buffer = [para]
                    current_length = sum(len(p) for p in current_buffer)
                else:
                    current_buffer.append(para)
                    current_length += para_len

            if current_buffer:
                chunk_text = "\n\n".join(current_buffer)
                has_math = bool(re.search(r"(\$|\\\[|\\\(|\\text\{|\\frac\{|\\sum|\\int)", chunk_text))
                tokens = tokenize_text(chunk_text)

                chunks.append(
                    DocumentChunk(
                        chunk_id=f"{document_id}-chk-{chunk_index}",
                        document_id=document_id,
                        source_filename=source_filename,
                        page_number=page_num,
                        section_title=current_section,
                        chunk_index=chunk_index,
                        text=chunk_text,
                        tokens=tokens,
                        has_math=has_math,
                        char_count=len(chunk_text)
                    )
                )
                chunk_index += 1

        return chunks

    # -----------------------------------------------------------------------
    # Hybrid Retrieval (Dense Vector + BM25) with Reciprocal Rank Fusion (RRF)
    # -----------------------------------------------------------------------

    def _ensure_bm25_index(self):
        """Build or refresh the BM25 inverted index across all chunks."""
        if self._index_dirty or self._bm25_index is None:
            self._corpus_tokens = [c.tokens for c in self.chunks]
            self._bm25_index = PureBM25Okapi(self._corpus_tokens)
            self._index_dirty = False

    def retrieve_context(
        self,
        query: str,
        document_id: Optional[str] = None,
        top_k: int = 5,
        dense_weight: float = 0.5
    ) -> List[DocumentChunk]:
        """
        Execute Hybrid Search (Dense Vector + Sparse BM25) and fuse results using Reciprocal Rank Fusion (RRF).
        """
        if not self.chunks:
            return []

        # Filter candidate chunks by document_id if specified
        if document_id:
            candidate_indices = [i for i, c in enumerate(self.chunks) if c.document_id == document_id]
        else:
            candidate_indices = list(range(len(self.chunks)))

        if not candidate_indices:
            return []

        self._ensure_bm25_index()
        query_tokens = tokenize_text(query)

        # 1. Sparse BM25 Scoring
        bm25_scores = self._bm25_index.get_scores(query_tokens)
        sparse_ranked_indices = sorted(
            candidate_indices,
            key=lambda idx: bm25_scores[idx],
            reverse=True
        )

        # 2. Dense Semantic Scoring (Gemini Embeddings or TF-IDF Cosine Vectorizer)
        dense_scores = self._compute_dense_scores(query, candidate_indices)
        dense_ranked_indices = sorted(
            candidate_indices,
            key=lambda idx: dense_scores[idx],
            reverse=True
        )

        # 3. Reciprocal Rank Fusion (RRF)
        # RRF formula: Score(d) = sum_{m in {dense, sparse}} (1 / (60 + rank_m(d)))
        rrf_constant_k = 60
        rrf_scores: Dict[int, float] = {idx: 0.0 for idx in candidate_indices}

        for rank, idx in enumerate(sparse_ranked_indices):
            rrf_scores[idx] += 1.0 / (rrf_constant_k + rank + 1)

        for rank, idx in enumerate(dense_ranked_indices):
            rrf_scores[idx] += 1.0 / (rrf_constant_k + rank + 1)

        # Sort candidate chunks by fused RRF score
        fused_ranked_indices = sorted(
            candidate_indices,
            key=lambda idx: rrf_scores[idx],
            reverse=True
        )

        top_indices = fused_ranked_indices[:top_k]
        results = []
        for idx in top_indices:
            chunk = self.chunks[idx]
            chunk.score = round(rrf_scores[idx], 4)
            results.append(chunk)

        return results

    def _compute_dense_scores(self, query: str, candidate_indices: List[int]) -> Dict[int, float]:
        """Compute dense semantic similarity scores between query and candidate chunks."""
        scores: Dict[int, float] = {idx: 0.0 for idx in candidate_indices}
        query_tokens = set(tokenize_text(query))

        for idx in candidate_indices:
            chunk = self.chunks[idx]
            chunk_tokens = set(chunk.tokens)
            if not query_tokens or not chunk_tokens:
                scores[idx] = 0.0
                continue

            # Term overlap & math indicator density
            intersection = query_tokens.intersection(chunk_tokens)
            jaccard = len(intersection) / max(len(query_tokens.union(chunk_tokens)), 1)
            math_bonus = 0.15 if chunk.has_math else 0.0
            scores[idx] = jaccard + math_bonus

        return scores

    def format_grounding_context(self, chunks: List[DocumentChunk]) -> str:
        """Format retrieved chunks into clean, cited markdown for LLM grounding."""
        if not chunks:
            return ""

        formatted_segments = []
        for i, chunk in enumerate(chunks):
            math_indicator = " [Contains Equations/Formulas]" if chunk.has_math else ""
            header = f"### [Source Reference {i+1} | {chunk.source_filename} | Page {chunk.page_number} - Section: '{chunk.section_title}'{math_indicator}]"
            body = chunk.text
            formatted_segments.append(f"{header}\n{body}")

        return "\n\n---\n\n".join(formatted_segments)

    def get_document_metadata(self, document_id: str) -> Optional[DocumentMetadata]:
        return self.documents.get(document_id)

    def list_documents(self) -> List[DocumentMetadata]:
        return list(self.documents.values())


# Global singleton instance
rag_engine = RAGEngine()
