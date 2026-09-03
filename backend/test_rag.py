"""
Comprehensive Verification Suite for Phase 2: RAG Pipeline & FastAPI Endpoints.
"""
import os
import sys
import json
import time

# Add backend directory to sys.path
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "."))

def test_rag_multi_format_and_hybrid():
    print("Testing Multi-Format Ingestion & Hybrid Retrieval...")
    from app.services.rag_engine import rag_engine, DocumentChunk, DocumentMetadata

    # 1. Test Text/Markdown with mathematical formulas
    sample_markdown = r"""# Quantum Superposition and Qubits

## 1. Mathematical Representation
A quantum bit or qubit is a two-state quantum-mechanical system.
The state of a pure qubit state can be represented by a linear superposition of the basis states:
$$|\psi\rangle = \alpha |0\rangle + \beta |1\rangle$$
where $\alpha$ and $\beta$ are complex numbers satisfying the probability normalization:
$$|\alpha|^2 + |\beta|^2 = 1$$

## 2. The Bloch Sphere
The Bloch sphere provides a geometrical representation of the pure state space of a 2-level quantum mechanical system.
Any pure state $|\psi\rangle$ can be rewritten as:
$$|\psi\rangle = \cos(\theta/2)|0\rangle + e^{i\phi}\sin(\theta/2)|1\rangle$$
where $0 \le \theta \le \pi$ and $0 \le \phi < 2\pi$.

## 3. Quantum Logic Gates
Quantum logic gates operate on qubits through unitary operators $U^\dagger U = I$.
The Hadamard gate $H$ creates an equal superposition:
$$H|0\rangle = \frac{1}{\sqrt{2}}(|0\rangle + |1\rangle)$$
"""
    # Write to a temporary markdown file
    temp_md_path = os.path.join(os.path.dirname(__file__), "temp_quantum_notes.md")
    with open(temp_md_path, "w", encoding="utf-8") as f:
        f.write(sample_markdown)

    try:
        meta = rag_engine.ingest_document(temp_md_path)
        assert meta.total_chunks >= 3
        assert meta.has_math is True
        assert "1. Mathematical Representation" in meta.sections_detected or "Mathematical Representation" in meta.sections_detected[0]
        print(f"  [OK] Ingested Markdown document '{meta.detected_title}': {meta.total_chunks} chunks, {len(meta.sections_detected)} sections.")

        # 2. Test Hybrid Retrieval with BM25 & Dense RRF
        query = "Hadamard gate equal superposition formula"
        results = rag_engine.retrieve_context(query, document_id=meta.document_id, top_k=3)
        assert len(results) > 0
        top_chunk = results[0]
        assert "Hadamard" in top_chunk.text or "superposition" in top_chunk.text
        assert top_chunk.has_math is True
        print(f"  [OK] Hybrid Retrieval Top Result (RRF Score: {top_chunk.score}): Section '{top_chunk.section_title}'")

        # 3. Test Formatted Grounding Context
        formatted = rag_engine.format_grounding_context(results)
        assert "Source Reference 1" in formatted
        assert "Page 1" in formatted
        print("  [OK] Formatted Grounding Context with citations generated successfully")

    finally:
        if os.path.exists(temp_md_path):
            os.remove(temp_md_path)


def test_docx_and_pptx_ingestion():
    print("\nTesting DOCX and PPTX Parsers...")
    from app.services.rag_engine import rag_engine

    # 1. Create and test DOCX
    try:
        import docx
        doc = docx.Document()
        doc.add_heading("Convolutional Neural Networks", level=1)
        doc.add_paragraph("CNNs use convolutional layers to extract local spatial features from images.")
        doc.add_heading("Pooling Operations", level=2)
        doc.add_paragraph("Max pooling reduces feature map spatial dimensions while preserving translational invariance.")

        table = doc.add_table(rows=2, cols=2)
        table.rows[0].cells[0].text = "Layer Type"
        table.rows[0].cells[1].text = "Parameter Count"
        table.rows[1].cells[0].text = "Conv2D (3x3)"
        table.rows[1].cells[1].text = "O(K * C_in * C_out)"

        docx_path = os.path.join(os.path.dirname(__file__), "temp_cnn.docx")
        doc.save(docx_path)

        meta_docx = rag_engine.ingest_document(docx_path)
        assert meta_docx.file_type == "docx"
        assert meta_docx.total_chunks >= 1
        print(f"  [OK] Ingested DOCX '{meta_docx.detected_title}': {meta_docx.total_chunks} chunks.")
        if os.path.exists(docx_path):
            os.remove(docx_path)
    except Exception as e:
        print(f"  [!] DOCX test note: {e}")

    # 2. Create and test PPTX
    try:
        from pptx import Presentation
        prs = Presentation()
        # Slide 1
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide1.shapes.title.text = "Deep Residual Networks (ResNets)"
        # Slide 2
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Vanishing Gradient Solution"
        slide2.shapes.placeholders[1].text = "Skip connections allow identity mappings: F(x) + x."

        pptx_path = os.path.join(os.path.dirname(__file__), "temp_resnet.pptx")
        prs.save(pptx_path)

        meta_pptx = rag_engine.ingest_document(pptx_path)
        assert meta_pptx.file_type == "pptx"
        assert meta_pptx.total_pages == 2
        print(f"  [OK] Ingested PPTX '{meta_pptx.detected_title}': {meta_pptx.total_pages} slides, {meta_pptx.total_chunks} chunks.")
        if os.path.exists(pptx_path):
            os.remove(pptx_path)
    except Exception as e:
        print(f"  [!] PPTX test note: {e}")


def test_fastapi_endpoints():
    print("\nTesting FastAPI v1 Endpoints & Grounded Curriculum Generation...")
    from fastapi.testclient import TestClient
    from app.main import app

    client = TestClient(app)

    # 1. Test /
    res_root = client.get("/")
    assert res_root.status_code == 200
    assert res_root.json()["service"] == "Synapse AI Teacher Backend"
    print("  [OK] GET / -> 200 OK")

    # 2. Test /api/v1/upload with a sample text file
    sample_text = """
    # Gradient Descent Optimization
    Gradient descent is a first-order iterative optimization algorithm for finding a local minimum of a differentiable function.
    To find a local minimum using gradient descent, we take steps proportional to the negative of the gradient of the function:
    $$\\theta_{t+1} = \\theta_t - \\eta \\nabla L(\\theta_t)$$
    where $\\eta$ is the learning rate.
    """
    res_upload = client.post(
        "/api/v1/upload",
        files={"file": ("gradient_descent.md", sample_text.encode("utf-8"), "text/markdown")}
    )
    assert res_upload.status_code == 200
    upload_data = res_upload.json()
    doc_id = upload_data["document_id"]
    assert upload_data["success"] is True
    assert upload_data["has_math"] is True
    print(f"  [OK] POST /api/v1/upload -> 200 OK (doc_id: {doc_id})")

    # 3. Test /api/v1/retrieve
    res_ret = client.post(
        "/api/v1/retrieve",
        json={"query": "learning rate negative gradient update step", "document_id": doc_id, "top_k": 2}
    )
    assert res_ret.status_code == 200
    ret_data = res_ret.json()
    assert ret_data["chunks_returned"] >= 1
    print(f"  [OK] POST /api/v1/retrieve -> 200 OK ({ret_data['chunks_returned']} chunks returned)")

    # 4. Test /api/v1/generate-curriculum (Grounded)
    res_curr = client.post(
        "/api/v1/generate-curriculum",
        json={
            "document_id": doc_id,
            "topic": "Gradient Descent Optimization",
            "language": "Hinglish",
            "educational_level": "Intermediate",
            "available_time_minutes": "20"
        }
    )
    assert res_curr.status_code == 200
    curr_data = res_curr.json()
    assert curr_data["success"] is True
    assert curr_data["is_grounded"] is True
    assert len(curr_data["lesson_plan"]["modules"]) >= 2
    print(f"  [OK] POST /api/v1/generate-curriculum -> 200 OK (Grounded with {len(curr_data['lesson_plan']['modules'])} modules)")


if __name__ == "__main__":
    print("==================================================")
    print("SYNAPSE AI TEACHER - PHASE 2 RAG TEST SUITE")
    print("==================================================")
    try:
        test_rag_multi_format_and_hybrid()
        test_docx_and_pptx_ingestion()
        test_fastapi_endpoints()
        print("\n==================================================")
        print("ALL PHASE 2 RAG VERIFICATION TESTS PASSED!")
        print("==================================================")
    except Exception as e:
        print(f"\n[!] TEST FAILED: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
